"""Generate a white-theme presentation for Market Gap Finder."""

from __future__ import annotations

import sys
from pathlib import Path

import lightgbm as lgb
import matplotlib
import numpy as np
import pandas as pd

matplotlib.use("Agg")
import matplotlib.pyplot as plt
import japanize_matplotlib  # noqa: F401

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

PROJECT_ROOT = Path(__file__).resolve().parents[1]
if str(PROJECT_ROOT) not in sys.path:
    sys.path.insert(0, str(PROJECT_ROOT))

from src.analyze.ml_model import compute_shap_values, train_cv, compute_market_gap, prepare_features, DEFAULT_PARAMS, DEFAULT_NUM_ROUNDS

FONT_NAME = "Meiryo"
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

WHITE_BG = RGBColor(0xFA, 0xFA, 0xFC)
BLUE = RGBColor(0x2D, 0x5B, 0xE3)
BLUE_LIGHT = RGBColor(0xE8, 0xEE, 0xFD)
GREEN = RGBColor(0x0F, 0x9D, 0x58)
RED = RGBColor(0xDB, 0x44, 0x37)
GRAY_CARD = RGBColor(0xF0, 0xF2, 0xF5)
TEXT = RGBColor(0x22, 0x28, 0x33)
MUTED = RGBColor(0x66, 0x6F, 0x7A)
LINE = RGBColor(0xD7, 0xDE, 0xEA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

CHART_DIR = PROJECT_ROOT / "docs" / "eda_charts"
OUTPUT_PATH = PROJECT_ROOT / "docs" / "market_gap_finder_presentation.pptx"
TOKYO_INTEGRATED = PROJECT_ROOT / "data" / "processed" / "tokyo_integrated.csv"
MODEL_PATH = PROJECT_ROOT / "models" / "tokyo_lightgbm.txt"
DASHBOARD_PATH = CHART_DIR / "dashboard_screenshot.png"


def set_slide_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE_BG


def add_rect(slide, left, top, width, height, fill_color, line_color=LINE, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = line_color
    shape.line.width = Pt(1)
    if radius and hasattr(shape, "adjustments") and len(shape.adjustments) > 0:
        shape.adjustments[0] = 0.08
    return shape


def format_paragraph(paragraph, font_size, color, bold=False, align=PP_ALIGN.LEFT) -> None:
    paragraph.alignment = align
    paragraph.space_after = Pt(0)
    paragraph.space_before = Pt(0)
    for run in paragraph.runs:
        run.font.name = FONT_NAME
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    *,
    font_size=18,
    color=TEXT,
    bold=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
):
    textbox = slide.shapes.add_textbox(left, top, width, height)
    frame = textbox.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    format_paragraph(paragraph, font_size, color, bold=bold, align=align)
    return textbox


def add_paragraphs(
    slide,
    left,
    top,
    width,
    height,
    items,
    *,
    font_size=16,
    color=TEXT,
    bullet=False,
    bullet_color=BLUE,
    spacing=5,
):
    textbox = slide.shapes.add_textbox(left, top, width, height)
    frame = textbox.text_frame
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.space_after = Pt(spacing)
        if bullet:
            bullet_run = paragraph.add_run()
            bullet_run.text = "• "
            bullet_run.font.name = FONT_NAME
            bullet_run.font.size = Pt(font_size)
            bullet_run.font.color.rgb = bullet_color
            text_run = paragraph.add_run()
            text_run.text = item
            text_run.font.name = FONT_NAME
            text_run.font.size = Pt(font_size)
            text_run.font.color.rgb = color
        else:
            paragraph.text = item
            format_paragraph(paragraph, font_size, color)
    return textbox


def add_title(slide, index_label: str, title: str) -> None:
    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.55), Inches(0.38), Inches(0.12), Inches(0.52))
    accent.fill.solid()
    accent.fill.fore_color.rgb = BLUE
    accent.line.fill.background()
    add_textbox(slide, Inches(0.78), Inches(0.24), Inches(11.4), Inches(0.5), f"{index_label}: {title}", font_size=28, bold=True)


def add_card_title(slide, left, top, width, title, color=BLUE) -> None:
    add_textbox(slide, left, top, width, Inches(0.28), title, font_size=16, color=color, bold=True)


def add_table(
    slide,
    left,
    top,
    width,
    height,
    rows,
    *,
    font_size=12,
    header_fill=BLUE,
    col_widths=None,
):
    row_count = len(rows)
    col_count = len(rows[0])
    shape = slide.shapes.add_table(row_count, col_count, left, top, width, height)
    table = shape.table

    if col_widths is not None:
        for idx, col_width in enumerate(col_widths):
            table.columns[idx].width = col_width

    row_height = int(height / row_count)
    for row_idx in range(row_count):
        table.rows[row_idx].height = row_height

    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(value)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for paragraph in cell.text_frame.paragraphs:
                paragraph.word_wrap = True
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.name = FONT_NAME
                    run.font.size = Pt(font_size)
                    run.font.bold = row_idx == 0
                    run.font.color.rgb = WHITE if row_idx == 0 else TEXT
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_fill if row_idx == 0 else (BLUE_LIGHT if row_idx % 2 == 1 else WHITE)
    return shape


def add_picture_or_placeholder(slide, path: Path, left, top, width, height, placeholder_text: str) -> None:
    if path.exists():
        slide.shapes.add_picture(str(path), left, top, width=width, height=height)
        return
    add_rect(slide, left, top, width, height, GRAY_CARD)
    add_textbox(
        slide,
        left + Inches(0.2),
        top + height / 2 - Inches(0.22),
        width - Inches(0.4),
        Inches(0.5),
        placeholder_text,
        font_size=14,
        color=MUTED,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )


def setup_chart_style() -> None:
    plt.rcParams.update(
        {
            "font.family": FONT_NAME,
            "font.size": 11,
            "figure.facecolor": "#FAFAFC",
            "axes.facecolor": "#FFFFFF",
            "axes.edgecolor": "#D7DEEA",
            "axes.labelcolor": "#222833",
            "text.color": "#222833",
            "xtick.color": "#59616B",
            "ytick.color": "#59616B",
            "grid.color": "#E8EDF5",
            "grid.linewidth": 0.7,
            "grid.alpha": 0.9,
        }
    )


def generate_eda_charts(out_dir: Path | str = CHART_DIR) -> Path:
    out_dir = Path(out_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    setup_chart_style()

    tokyo_df = pd.read_csv(TOKYO_INTEGRATED)

    # Train Tokyo-only model and get OOF predictions for scatter plot
    cv_results = train_cv(tokyo_df)
    tokyo_gap_df = compute_market_gap(tokyo_df, cv_results["oof_predictions"])

    density_df = (
        tokyo_df.groupby("jis_mesh", as_index=False)
        .agg(total=("restaurant_count", "sum"), lat=("lat", "first"), lng=("lng", "first"))
        .dropna(subset=["lat", "lng"])
    )

    fig, ax = plt.subplots(figsize=(7.2, 5.2))
    points = ax.scatter(
        density_df["lng"],
        density_df["lat"],
        c=density_df["total"],
        cmap="YlOrRd",
        s=np.clip(density_df["total"], 1, 80) * 1.5 + 6,
        alpha=0.8,
        edgecolors="none",
    )
    colorbar = fig.colorbar(points, ax=ax, shrink=0.82, pad=0.02)
    colorbar.set_label("飲食店数")
    ax.set_title("東京の店舗集積マップ", fontsize=15, fontweight="bold")
    ax.set_xlabel("経度")
    ax.set_ylabel("緯度")
    ax.grid(True)
    for name, x, y, dx, dy in [
        ("新宿", 139.700, 35.690, -0.030, 0.010),
        ("渋谷", 139.702, 35.660, -0.030, -0.005),
        ("池袋", 139.711, 35.730, 0.010, 0.012),
        ("東京駅", 139.767, 35.681, 0.010, -0.008),
        ("上野", 139.774, 35.714, 0.010, 0.010),
    ]:
        ax.annotate(
            name,
            xy=(x, y),
            xytext=(x + dx, y + dy),
            fontsize=10,
            color="#222833",
            arrowprops={"arrowstyle": "->", "color": "#2D5BE3", "lw": 1.1},
            bbox={"boxstyle": "round,pad=0.2", "fc": "#FFFFFF", "ec": "#D7DEEA"},
        )
    fig.tight_layout()
    fig.savefig(out_dir / "tokyo_density_map.png", dpi=180)
    plt.close(fig)

    genre_df = (
        tokyo_df.groupby("unified_genre", as_index=False)["restaurant_count"]
        .sum()
        .sort_values("restaurant_count", ascending=True)
    )
    fig, ax = plt.subplots(figsize=(5.2, 4.3))
    colors = ["#2D5BE3"] * len(genre_df)
    if len(colors) > 0:
        colors[-1] = "#0F9D58"
    bars = ax.barh(genre_df["unified_genre"], genre_df["restaurant_count"], color=colors)
    ax.set_title("ジャンル別店舗数", fontsize=14, fontweight="bold")
    ax.set_xlabel("店舗数")
    ax.grid(True, axis="x")
    for bar, value in zip(bars, genre_df["restaurant_count"]):
        ax.text(value + genre_df["restaurant_count"].max() * 0.02, bar.get_y() + bar.get_height() / 2, f"{int(value):,}", va="center", fontsize=9)
    fig.tight_layout()
    fig.savefig(out_dir / "tokyo_genre_bar.png", dpi=180)
    plt.close(fig)

    counts = tokyo_df["restaurant_count"].fillna(0).to_numpy()
    fig, axes = plt.subplots(1, 2, figsize=(7.2, 3.6))
    axes[0].hist(counts, bins=40, color="#2D5BE3", edgecolor="white")
    axes[0].set_title("元データ")
    axes[0].set_xlabel("店舗数")
    axes[0].set_ylabel("メッシュ数")
    axes[0].grid(True)
    axes[1].hist(np.log1p(counts), bins=40, color="#0F9D58", edgecolor="white")
    axes[1].set_title("log変換後")
    axes[1].set_xlabel("log1p(店舗数)")
    axes[1].set_ylabel("メッシュ数")
    axes[1].grid(True)
    fig.suptitle("店舗数分布の偏りと対数変換", fontsize=13, fontweight="bold")
    fig.tight_layout()
    fig.savefig(out_dir / "tokyo_histogram.png", dpi=180, bbox_inches="tight")
    plt.close(fig)

    scatter_df = tokyo_gap_df.dropna(subset=["actual_log_count", "predicted_log_count"])
    actual_count = scatter_df["restaurant_count"].to_numpy()
    predicted_count = np.expm1(scatter_df["predicted_log_count"].to_numpy())
    rmse_count = np.sqrt(np.mean((actual_count - predicted_count) ** 2))

    # 全体散布図
    fig, ax = plt.subplots(figsize=(5.8, 4.8))
    ax.scatter(actual_count, predicted_count, color="#2D5BE3", alpha=0.35, s=18, edgecolors="none")
    upper = max(actual_count.max(), predicted_count.max())
    ax.plot([0, upper], [0, upper], linestyle="--", color="#DB4437", linewidth=1.5, label="y = x（完全一致）")
    ax.set_title("実際の店舗数 vs 予測店舗数（全体）", fontsize=14, fontweight="bold")
    ax.set_xlabel("実際の店舗数")
    ax.set_ylabel("予測店舗数")
    ax.text(0.05, 0.93, f"RMSE = {rmse_count:.1f}", transform=ax.transAxes, fontsize=12, bbox={"boxstyle": "round,pad=0.3", "fc": "#FFFFFF", "ec": "#D7DEEA"})
    ax.legend(loc="lower right", fontsize=10)
    ax.grid(True)
    fig.tight_layout()
    fig.savefig(out_dir / "tokyo_scatter_actual_vs_pred.png", dpi=180)
    plt.close(fig)

    # 0〜25拡大散布図
    mask = (actual_count <= 25) & (predicted_count <= 25)
    fig, ax = plt.subplots(figsize=(5.8, 4.8))
    ax.scatter(actual_count[mask], predicted_count[mask], color="#2D5BE3", alpha=0.35, s=18, edgecolors="none")
    ax.plot([0, 25], [0, 25], linestyle="--", color="#DB4437", linewidth=1.5, label="y = x（完全一致）")
    ax.set_title("実際の店舗数 vs 予測店舗数（0〜25店に拡大）", fontsize=14, fontweight="bold")
    ax.set_xlabel("実際の店舗数")
    ax.set_ylabel("予測店舗数")
    ax.set_xlim(0, 25)
    ax.set_ylim(0, 25)
    rmse_25 = np.sqrt(np.mean((actual_count[mask] - predicted_count[mask]) ** 2))
    ax.text(0.05, 0.93, f"RMSE = {rmse_25:.1f}（0〜25店）", transform=ax.transAxes, fontsize=12, bbox={"boxstyle": "round,pad=0.3", "fc": "#FFFFFF", "ec": "#D7DEEA"})
    ax.legend(loc="lower right", fontsize=10)
    ax.grid(True)
    fig.tight_layout()
    fig.savefig(out_dir / "tokyo_scatter_zoom25.png", dpi=180)
    plt.close(fig)

    # Train fresh Tokyo model (13 features) for SHAP
    features_df, target = prepare_features(tokyo_df)
    cat_features = ["genre_encoded"] if "genre_encoded" in features_df.columns else []
    train_data = lgb.Dataset(features_df, label=target, categorical_feature=cat_features)
    tokyo_model = lgb.train(DEFAULT_PARAMS, train_data, num_boost_round=DEFAULT_NUM_ROUNDS,
                            valid_sets=[train_data], callbacks=[lgb.log_evaluation(0)])
    shap_values, features = compute_shap_values(tokyo_model, tokyo_df)
    shap_array = shap_values[0] if isinstance(shap_values, list) else shap_values
    importance = pd.DataFrame({"feature": features.columns, "importance": np.abs(shap_array).mean(axis=0)})
    top10 = importance.sort_values("importance", ascending=False).head(10).sort_values("importance", ascending=True)
    fig, ax = plt.subplots(figsize=(6.2, 4.8))
    ax.barh(top10["feature"], top10["importance"], color="#2D5BE3")
    ax.set_title("SHAP重要度 Top10", fontsize=14, fontweight="bold")
    ax.set_xlabel("平均|SHAP値|")
    ax.grid(True, axis="x")
    fig.tight_layout()
    fig.savefig(out_dir / "shap_importance_top10.png", dpi=180)
    plt.close(fig)

    # 人口データの欠損率（4分の1メッシュ単位、ゼロ値＝e-Stat未提供）
    mesh_df = tokyo_df.drop_duplicates(subset="jis_mesh")
    n_meshes = len(mesh_df)
    pop_cols_labels = {
        "population": "総人口",
        "pop_working": "労働人口",
        "pop_adult": "成人",
        "pop_elderly": "高齢者",
        "households": "世帯数",
        "single_households": "単身世帯",
        "young_single": "若年単身",
    }
    zero_rates = []
    labels = []
    for col, label in pop_cols_labels.items():
        if col in mesh_df.columns:
            rate = (mesh_df[col] == 0).sum() / n_meshes * 100
            zero_rates.append(rate)
            labels.append(label)
    fig, ax = plt.subplots(figsize=(5.2, 3.2))
    bars = ax.barh(labels, zero_rates, color="#2D5BE3")
    ax.set_title(f"人口データのゼロ値率（メッシュ単位, n={n_meshes:,}）", fontsize=13, fontweight="bold")
    ax.set_xlabel("ゼロ値の割合 (%)")
    ax.set_xlim(0, max(zero_rates) * 1.3 if zero_rates else 5)
    ax.grid(True, axis="x")
    for bar, val in zip(bars, zero_rates):
        ax.text(val + 0.1, bar.get_y() + bar.get_height() / 2, f"{val:.1f}%", va="center", fontsize=9)
    fig.tight_layout()
    fig.savefig(out_dir / "tokyo_pop_missing.png", dpi=180)
    plt.close(fig)

    return out_dir


def build_cover_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_rect(slide, Inches(0.7), Inches(0.9), Inches(11.9), Inches(5.7), WHITE, line_color=LINE)
    add_rect(slide, Inches(0.95), Inches(1.18), Inches(2.0), Inches(0.18), BLUE, line_color=BLUE, radius=False)
    add_textbox(slide, Inches(0.95), Inches(1.55), Inches(11), Inches(0.9), "市場空白地帯発見システム", font_size=30, bold=True)
    add_textbox(slide, Inches(0.95), Inches(2.3), Inches(11), Inches(0.45), "Market Gap Finder", font_size=20, color=BLUE, bold=True)
    add_textbox(slide, Inches(0.95), Inches(3.15), Inches(11), Inches(0.7), "「どこに出店すればいいか？」をデータで答える", font_size=24, color=TEXT, bold=True)
    add_textbox(slide, Inches(0.95), Inches(4.45), Inches(8), Inches(0.35), "Vibe Coding with Claude Code + Codex MCP", font_size=14, color=MUTED)
    add_textbox(slide, Inches(0.95), Inches(4.82), Inches(8), Inches(0.35), "2026-04-08", font_size=14, color=MUTED)


def build_problem_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_title(slide, "1", "ビジネス課題")

    add_rect(slide, Inches(0.65), Inches(1.0), Inches(12.0), Inches(1.35), BLUE_LIGHT)
    add_textbox(slide, Inches(0.95), Inches(1.18), Inches(11.4), Inches(0.95), "「なぜ駅前にはラーメン屋が多いのか？」「この街にカフェが少ないのはなぜ？」── こうした疑問をデータで説明できたら面白い。それがこのシステムの出発点。", font_size=19, bold=True)

    add_rect(slide, Inches(0.65), Inches(2.6), Inches(12.0), Inches(1.4), GRAY_CARD)
    add_card_title(slide, Inches(0.95), Inches(2.78), Inches(3.0), "課題", color=RED)
    add_textbox(slide, Inches(0.95), Inches(3.08), Inches(11.1), Inches(0.7), "飲食店テナントは頻繁に入れ替わる。出店判断の基準は各社・各担当で異なり、新しいエリアやジャンルへの横展開が難しい。", font_size=16)

    add_rect(slide, Inches(0.65), Inches(4.2), Inches(12.0), Inches(1.4), BLUE_LIGHT)
    add_card_title(slide, Inches(0.95), Inches(4.38), Inches(3.0), "解決アプローチ", color=GREEN)
    add_textbox(slide, Inches(0.95), Inches(4.68), Inches(11.1), Inches(0.7), "公開データから「どこに、どんなジャンルの店が足りないか？」を数値化し、エリアやジャンルを問わず使える出店基準をつくる。", font_size=16)

    add_textbox(slide, Inches(0.65), Inches(5.9), Inches(4.5), Inches(0.35), "このテーマを選んだ理由", font_size=17, bold=True)
    add_paragraphs(
        slide,
        Inches(0.85),
        Inches(6.2),
        Inches(11.4),
        Inches(1.0),
        [
            "頻繁に入れ替わる飲食テナントに対して、データでニーズを可視化したかった",
            "各社バラバラの出店基準を、データから導く普遍的な基準に置き換えたかった",
            "講座で学んだML・データ分析・可視化を統合的に活かせるテーマ",
        ],
        font_size=14,
        color=MUTED,
        bullet=True,
    )


def build_data_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_title(slide, "2", "データの紹介")

    add_table(
        slide,
        Inches(0.55),
        Inches(0.95),
        Inches(12.2),
        Inches(1.62),
        [
            ["カテゴリ", "知りたいこと", "データソース"],
            ["需要（お客さんはいる？）", "その街にどれだけ人がいるか", "e-Stat（国勢調査）"],
            ["供給（ライバルは？）", "今どれだけ店があるか、評判は？", "ホットペッパー / Google Places"],
            ["立地条件（場所の特性は？）", "駅に近い？地価は？", "HeartRails / 国土数値情報"],
        ],
        font_size=12,
        col_widths=[Inches(2.9), Inches(5.2), Inches(4.1)],
    )

    add_table(
        slide,
        Inches(0.55),
        Inches(2.8),
        Inches(12.2),
        Inches(2.3),
        [
            ["API / データ", "役割", "粒度", "補足"],
            ["ホットペッパー", "店舗名・ジャンル・位置情報の取得", "店舗単位", "供給の主データ"],
            ["Google Places", "評価・レビュー数・補完的な店舗情報", "店舗単位", "評判の指標"],
            ["e-Stat", "人口・世帯・年齢構成", "250mメッシュ", "需要の指標"],
            ["HeartRails", "最寄駅と駅距離", "地点単位", "アクセス指標"],
            ["駅乗降客数", "駅の人流規模", "駅単位", "流動人口の近似"],
            ["国土数値情報", "地価データ", "地点単位", "立地の格を表現"],
        ],
        font_size=11,
        col_widths=[Inches(2.4), Inches(4.6), Inches(1.8), Inches(3.4)],
    )

    add_rect(slide, Inches(0.55), Inches(5.45), Inches(5.9), Inches(1.2), BLUE_LIGHT)
    add_card_title(slide, Inches(0.85), Inches(5.62), Inches(2.0), "対象エリア", color=BLUE)
    add_textbox(slide, Inches(0.85), Inches(5.95), Inches(5.3), Inches(0.45), "東京", font_size=18, bold=True)

    add_rect(slide, Inches(6.75), Inches(5.45), Inches(6.0), Inches(1.2), GRAY_CARD)
    add_card_title(slide, Inches(7.05), Inches(5.62), Inches(2.0), "対象ジャンル", color=BLUE)
    add_textbox(slide, Inches(7.05), Inches(5.95), Inches(5.5), Inches(0.45), "カフェ / ラーメン / 居酒屋 / 焼肉 / 和食 / 中華 / イタリアン / その他", font_size=14, bold=True)


def build_mesh_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_title(slide, "2.5", "分析の単位 ── メッシュとは？")

    add_rect(slide, Inches(0.65), Inches(1.0), Inches(12.0), Inches(1.0), BLUE_LIGHT)
    add_textbox(slide, Inches(0.95), Inches(1.25), Inches(11.4), Inches(0.45), "地図の上に250m四方のマス目（格子）を敷いて、マスごとに「飲食店は何軒？」「人口は？」を集計する。", font_size=18, bold=True)

    add_rect(slide, Inches(0.65), Inches(2.25), Inches(5.65), Inches(3.7), GRAY_CARD)
    add_card_title(slide, Inches(0.95), Inches(2.45), Inches(2.5), "ポイント")
    add_paragraphs(
        slide,
        Inches(0.95),
        Inches(2.8),
        Inches(4.95),
        Inches(2.7),
        [
            "1マス = 250m × 250m（だいたいコンビニ2〜3軒分の距離感）",
            "東京で合計 1,233マス",
            "マスごとに「需要・供給・立地」のデータを紐づけて比較",
        ],
        font_size=15,
        bullet=True,
    )

    grid = slide.shapes.add_table(4, 4, Inches(7.0), Inches(2.45), Inches(4.9), Inches(3.2)).table
    for row in range(4):
        for col in range(4):
            cell = grid.cell(row, col)
            cell.fill.solid()
            cell.fill.fore_color.rgb = BLUE_LIGHT if (row + col) % 2 == 0 else WHITE
    highlight = grid.cell(1, 1)
    highlight.fill.fore_color.rgb = BLUE
    highlight.text = "1メッシュ"
    for paragraph in highlight.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.name = FONT_NAME
            run.font.size = Pt(16)
            run.font.bold = True
            run.font.color.rgb = WHITE

    add_rect(slide, Inches(0.65), Inches(6.1), Inches(12.0), Inches(0.95), BLUE_LIGHT)
    add_textbox(slide, Inches(0.95), Inches(6.35), Inches(11.2), Inches(0.4), "「渋谷区」のような大きい単位だと駅前と住宅街の違いが見えない。250m単位なら「この通りの向こう側」レベルで出店チャンスを見つけられる。", font_size=16)


def build_eda_map_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_title(slide, "3a", "EDA ― 店舗集積マップ・ジャンル分布")
    add_picture_or_placeholder(slide, CHART_DIR / "tokyo_density_map.png", Inches(0.55), Inches(1.0), Inches(7.2), Inches(5.6), "密度マップ画像が見つかりません")
    add_picture_or_placeholder(slide, CHART_DIR / "tokyo_genre_bar.png", Inches(8.0), Inches(1.0), Inches(4.75), Inches(3.1), "ジャンル別棒グラフが見つかりません")
    add_rect(slide, Inches(8.0), Inches(4.35), Inches(4.75), Inches(2.25), GRAY_CARD)
    add_card_title(slide, Inches(8.25), Inches(4.55), Inches(2.0), "ここから言えること", color=RED)
    add_paragraphs(
        slide,
        Inches(8.25),
        Inches(4.95),
        Inches(4.2),
        Inches(1.4),
        [
            "主張1: 飲食店は一部エリアに極端に集中している",
            "主張2: ジャンルによって店舗数の差が10倍以上",
        ],
        font_size=15,
        bullet=True,
        bullet_color=RED,
    )


def build_eda_hist_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_title(slide, "3b", "EDA ― わかったこと → モデルへの反映")

    # 左上: ヒストグラム
    add_picture_or_placeholder(slide, CHART_DIR / "tokyo_histogram.png", Inches(0.55), Inches(0.95), Inches(5.2), Inches(2.9), "ヒストグラム画像が見つかりません")

    # 右上: 人口欠損率チャート
    add_picture_or_placeholder(slide, CHART_DIR / "tokyo_pop_missing.png", Inches(5.95), Inches(0.95), Inches(6.75), Inches(2.9), "人口欠損率チャートが見つかりません")

    # 左下: 気づき
    add_rect(slide, Inches(0.55), Inches(4.1), Inches(5.9), Inches(2.6), BLUE_LIGHT)
    add_card_title(slide, Inches(0.82), Inches(4.32), Inches(1.8), "気づき")
    add_paragraphs(
        slide,
        Inches(0.82),
        Inches(4.7),
        Inches(5.25),
        Inches(1.7),
        [
            "ほとんどのエリアは店舗1〜2軒、一部だけ100軒超 → データが偏っている",
            "「他ジャンルの店が多いエリア」ほど飲食店が多い → 商業集積が鍵",
            "人口データの欠損は0.4%（5/1,233メッシュ）→ ほぼ完全に揃っている",
        ],
        font_size=14,
        bullet=True,
    )

    # 右下: だからこうした
    add_rect(slide, Inches(6.6), Inches(4.1), Inches(6.1), Inches(2.6), GRAY_CARD)
    add_card_title(slide, Inches(6.87), Inches(4.32), Inches(2.0), "だからこうした", color=GREEN)
    add_paragraphs(
        slide,
        Inches(6.87),
        Inches(4.7),
        Inches(5.45),
        Inches(1.7),
        [
            "偏ったデータ → 対数変換（log）で公平に比較",
            "人口データの欠損0.4% → ゼロ埋めしてMLモデルに含む（除外しても精度差なし）",
        ],
        font_size=14,
        bullet=True,
        bullet_color=GREEN,
    )


def build_approach_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_title(slide, "4", "アプローチ ― MLで市場ギャップを予測")

    # 目的変数の説明
    add_rect(slide, Inches(0.55), Inches(1.0), Inches(12.2), Inches(1.6), BLUE_LIGHT)
    add_card_title(slide, Inches(0.85), Inches(1.18), Inches(3.0), "目的変数（何を予測するか）")
    add_paragraphs(
        slide,
        Inches(0.85),
        Inches(1.55),
        Inches(11.5),
        Inches(0.8),
        [
            "log1p(restaurant_count) ＝ 各メッシュ×ジャンルの飲食店数を対数変換した値",
            "対数変換する理由: 店舗数は0〜100超まで偏りが大きく、そのままでは少数の密集エリアに引きずられるため",
        ],
        font_size=15,
        bullet=True,
    )

    # 左カード: 特徴量（入力）
    add_rect(slide, Inches(0.55), Inches(2.85), Inches(5.95), Inches(2.3), GRAY_CARD)
    add_card_title(slide, Inches(0.85), Inches(3.05), Inches(3.4), "説明変数（26特徴量）", color=GREEN)
    add_paragraphs(
        slide,
        Inches(0.85),
        Inches(3.45),
        Inches(5.25),
        Inches(1.4),
        [
            "人口・世帯構成（11）、競合環境（4）、駅アクセス（3）",
            "地価・Google評価（3）、ジャンル関連（5）",
        ],
        font_size=15,
        bullet=True,
        bullet_color=GREEN,
    )

    # 右カード: 予測→ギャップ
    add_rect(slide, Inches(6.8), Inches(2.85), Inches(5.95), Inches(2.3), BLUE_LIGHT)
    add_card_title(slide, Inches(7.1), Inches(3.05), Inches(3.2), "市場ギャップの算出")
    add_paragraphs(
        slide,
        Inches(7.1),
        Inches(3.45),
        Inches(5.25),
        Inches(1.4),
        [
            "LightGBMで「本来何軒あるべきか」を予測",
            "予測値 − 実測値 = 市場ギャップ（大きい＝出店チャンス）",
        ],
        font_size=15,
        bullet=True,
    )

    # 下部: スコアバージョン比較テーブル
    add_table(
        slide,
        Inches(0.85),
        Inches(5.45),
        Inches(11.5),
        Inches(1.45),
        [
            ["", "ML予測のみ（デフォルト）", "ルールベース", "両方の組み合わせ"],
            ["方式", "ML予測ギャップをrank正規化", "6要素の重み付け合算", "ルール×0.6 + ML×0.4"],
            ["特徴", "精度が高い・シンプル", "解釈しやすい", "両方の長所"],
        ],
        font_size=11,
        col_widths=[Inches(1.2), Inches(3.8), Inches(3.2), Inches(3.3)],
    )


def build_feature_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_title(slide, "5", "モデルに入れたデータ（説明変数）")
    add_textbox(slide, Inches(0.65), Inches(0.95), Inches(12.0), Inches(0.4), "特徴量 = 「予測のヒントになるデータ」。人間が出店判断で見る情報を数値化したもの。", font_size=17, bold=True)
    add_table(
        slide,
        Inches(0.55),
        Inches(1.5),
        Inches(12.2),
        Inches(4.7),
        [
            ["カテゴリ", "個数", "具体例", "何を表すか"],
            ["人口", "7", "総人口、労働人口、成人、高齢者、世帯数、単身世帯、若年単身", "お客さんの数"],
            ["人口比率", "4", "労働人口率、高齢者率、単身率、若年単身率", "街の住民構成"],
            ["競合環境", "4", "他ジャンル店舗数、商業密度ランク、近隣平均店舗数、飽和度", "ライバルの多さ"],
            ["駅アクセス", "3", "最寄駅距離、乗降客数、アクセス指数", "人の流れ"],
            ["地価・評価", "3", "地価、Google平均評価、店あたりレビュー数", "エリアの格"],
            ["ジャンル関連", "5", "ジャンル種別、多様性、集中度(HHI)、近隣平均人口、Google密度", "飲食バリエーション"],
        ],
        font_size=11,
        col_widths=[Inches(1.8), Inches(0.8), Inches(7.0), Inches(2.6)],
    )
    add_textbox(slide, Inches(0.65), Inches(6.45), Inches(12.0), Inches(0.35), "補足: MLは26特徴量を使用（人口系含む）。人口欠損0.4%はゼロ埋め。Optuna 100 trialでチューニング済み。", font_size=14, color=MUTED)


def build_accuracy_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_title(slide, "6", "モデルの精度 ― 5-Fold クロスバリデーション")

    # 散布図（全体）
    add_picture_or_placeholder(slide, CHART_DIR / "tokyo_scatter_actual_vs_pred.png", Inches(0.55), Inches(0.95), Inches(4.2), Inches(3.4), "散布図画像が見つかりません")
    # 散布図（0〜25拡大）
    add_picture_or_placeholder(slide, CHART_DIR / "tokyo_scatter_zoom25.png", Inches(4.85), Inches(0.95), Inches(4.2), Inches(3.4), "拡大散布図が見つかりません")

    # なぜこの検証？
    add_rect(slide, Inches(9.2), Inches(0.95), Inches(3.5), Inches(3.4), BLUE_LIGHT)
    add_card_title(slide, Inches(9.4), Inches(1.1), Inches(3.1), "なぜこの検証方法？")
    add_paragraphs(
        slide,
        Inches(9.4),
        Inches(1.45),
        Inches(3.1),
        Inches(2.7),
        [
            "新しい時点のデータで精度を測るのが理想だが、時系列データが取れない",
            "学習データで評価するとリーケージになるため、店舗数分布が均等になるよう層化分割して汎化性能を測る",
        ],
        font_size=12,
        bullet=True,
    )

    # 具体的な方法
    add_rect(slide, Inches(0.55), Inches(4.55), Inches(6.0), Inches(2.2), GRAY_CARD)
    add_card_title(slide, Inches(0.82), Inches(4.72), Inches(3.0), "具体的な方法", color=GREEN)
    add_paragraphs(
        slide,
        Inches(0.82),
        Inches(5.05),
        Inches(5.4),
        Inches(1.5),
        [
            "東京の4,520行を店舗数の偏りが出ないよう5グループに層化分割",
            "StratifiedKFoldで店舗数ビンが均等になるよう分割",
            "4グループで学習→残り1グループで予測、を5回繰り返す",
        ],
        font_size=13,
        bullet=True,
        bullet_color=GREEN,
    )

    # 結果
    add_rect(slide, Inches(6.75), Inches(4.55), Inches(5.95), Inches(2.2), BLUE_LIGHT)
    add_card_title(slide, Inches(7.02), Inches(4.72), Inches(3.0), "結果")
    add_paragraphs(
        slide,
        Inches(7.02),
        Inches(5.05),
        Inches(5.4),
        Inches(1.5),
        [
            "RMSE = 5.0店（実店舗数スケール）＝ 平均で約5店の誤差",
            "中央値2店・平均4店のデータに対し、MAE = 1.6店（平均絶対誤差）",
        ],
        font_size=13,
        bullet=True,
    )


def build_shap_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_title(slide, "7", "何が予測に効いているか（SHAP）")
    add_picture_or_placeholder(slide, CHART_DIR / "shap_importance_top10.png", Inches(0.65), Inches(1.0), Inches(7.1), Inches(5.7), "SHAP画像が見つかりません")
    add_rect(slide, Inches(8.0), Inches(1.0), Inches(4.7), Inches(5.7), GRAY_CARD)
    add_card_title(slide, Inches(8.25), Inches(1.28), Inches(2.2), "解釈", color=RED)
    add_paragraphs(
        slide,
        Inches(8.25),
        Inches(1.75),
        Inches(4.1),
        Inches(2.0),
        [
            "ジャンル種別（SHAP=0.196）と他ジャンル店舗数（0.109）が上位",
            "人口系特徴量は含めているが重要度は低い（下位に位置）",
            "→ 「人口が多い」より「商業集積がある」方が飲食店数を決める",
        ],
        font_size=15,
        bullet=True,
        bullet_color=RED,
    )


def build_result_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_title(slide, "8", "課題点とうまくいった点")
    add_rect(slide, Inches(0.7), Inches(1.0), Inches(5.85), Inches(5.8), BLUE_LIGHT)
    add_card_title(slide, Inches(0.98), Inches(1.25), Inches(2.8), "うまくいった点", color=GREEN)
    add_paragraphs(
        slide,
        Inches(0.98),
        Inches(1.7),
        Inches(5.2),
        Inches(4.7),
        [
            "RMSE = 5.0店: 実店舗数スケールで平均約5店の誤差",
            "28特徴量（人口含む）でOptuna 100 trialチューニング",
            "ML予測のみでシンプルかつ高精度なスコアリング",
        ],
        font_size=16,
        bullet=True,
        bullet_color=GREEN,
    )
    add_rect(slide, Inches(6.75), Inches(1.0), Inches(5.85), Inches(5.8), GRAY_CARD)
    add_card_title(slide, Inches(7.03), Inches(1.25), Inches(2.8), "課題点", color=RED)
    add_paragraphs(
        slide,
        Inches(7.03),
        Inches(1.7),
        Inches(5.15),
        Inches(4.7),
        [
            "約半数のエリアが店舗数1軒（データの偏り）",
            "バックテスト未実施",
        ],
        font_size=16,
        bullet=True,
        bullet_color=RED,
    )


def build_business_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_title(slide, "9", "ビジネスインパクト")
    add_rect(slide, Inches(0.65), Inches(0.95), Inches(12.0), Inches(0.95), BLUE_LIGHT)
    add_textbox(slide, Inches(0.95), Inches(1.23), Inches(11.3), Inches(0.32), "「どの街に、どのジャンルの店が足りないか」を自動で見つけて、根拠付きで推薦する", font_size=19, bold=True)

    add_rect(slide, Inches(0.65), Inches(2.15), Inches(5.85), Inches(4.5), BLUE_LIGHT)
    add_card_title(slide, Inches(0.92), Inches(2.38), Inches(3.5), "新規エリア進出の判断")
    add_paragraphs(
        slide,
        Inches(0.92),
        Inches(2.78),
        Inches(5.2),
        Inches(3.5),
        [
            "東京都内の未出店エリアを、既存データだけで事前評価",
            "東京1,233メッシュ × 8ジャンル = 約9,800パターンの候補を一括で洗い出し",
        ],
        font_size=15,
        bullet=True,
    )

    add_rect(slide, Inches(6.75), Inches(2.15), Inches(5.85), Inches(4.5), GRAY_CARD)
    add_card_title(slide, Inches(7.02), Inches(2.38), Inches(3.5), "ジャンル戦略の意思決定", color=GREEN)
    add_paragraphs(
        slide,
        Inches(7.02),
        Inches(2.78),
        Inches(5.2),
        Inches(3.5),
        [
            "「このエリアには焼肉よりカフェの方がチャンスがある」をデータで裏付け",
            "ジャンルごとの市場ギャップを比較し、参入優先度を決定",
            "250mメッシュ粒度で「この通りの向こう側」レベルの精度",
        ],
        font_size=15,
        bullet=True,
        bullet_color=GREEN,
    )


def build_vibe_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_title(slide, "10", "Vibe Codingの所感")
    add_rect(slide, Inches(0.65), Inches(0.95), Inches(12.0), Inches(1.15), BLUE_LIGHT)
    add_textbox(slide, Inches(0.92), Inches(1.18), Inches(11.4), Inches(0.72), "自分の知識だけではここまでのモデルは作れなかった。AIが知識を補完してくれたからこそ、作りたいシステムを形にできた。ただしAIの提案を鵜呑みにしない判断力が不可欠。", font_size=17, bold=True)

    add_rect(slide, Inches(0.65), Inches(2.45), Inches(5.9), Inches(3.8), BLUE_LIGHT)
    add_card_title(slide, Inches(0.95), Inches(2.68), Inches(2.5), "AIに助けられた点", color=GREEN)
    add_paragraphs(
        slide,
        Inches(0.95),
        Inches(3.08),
        Inches(5.2),
        Inches(2.7),
        [
            "6種API + 26特徴量ML + 5エリア対応ダッシュボードを一人で構築",
            "知らなかった手法をAIが提案→知識の補完",
            "設計→実装→テストのサイクルが速く回せた",
        ],
        font_size=16,
        bullet=True,
        bullet_color=GREEN,
    )

    add_rect(slide, Inches(6.75), Inches(2.45), Inches(5.9), Inches(3.8), GRAY_CARD)
    add_card_title(slide, Inches(7.05), Inches(2.68), Inches(3.0), "人間が判断すべきだった点", color=RED)
    add_paragraphs(
        slide,
        Inches(7.05),
        Inches(3.08),
        Inches(5.15),
        Inches(2.7),
        [
            "AIが勝手に仕様変更を進めることがある",
            "R²やCIの解釈はドメイン知識が必要",
            "長い会話で過去の決定を忘れる",
        ],
        font_size=16,
        bullet=True,
        bullet_color=RED,
    )


def build_summary_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_title(slide, "11", "まとめ")
    add_table(
        slide,
        Inches(0.8),
        Inches(1.15),
        Inches(11.8),
        Inches(4.9),
        [
            ["項目", "内容"],
            ["課題", "出店判断の基準をデータで統一・定量化"],
            ["データ", "6種の公開データ × 東京 × 8ジャンル"],
            ["EDA", "商業集積が飲食店数を決める要因と判明"],
            ["手法", "LightGBM による市場ギャップ予測"],
            ["性能", "5-Fold CV RMSE=5.0店（東京エリア、層化分割、Optunaチューニング済み）"],
            ["活用", "東京都内の出店候補・ジャンル戦略をデータで裏付け"],
            ["開発", "AIと人間の判断力の掛け算で実現"],
        ],
        font_size=13,
        col_widths=[Inches(2.0), Inches(9.8)],
    )
    add_textbox(slide, Inches(4.1), Inches(6.35), Inches(5.2), Inches(0.45), "Thank you", font_size=26, color=BLUE, bold=True, align=PP_ALIGN.CENTER)


def build_presentation() -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    build_cover_slide(prs)
    build_problem_slide(prs)
    build_data_slide(prs)
    build_mesh_slide(prs)
    build_eda_map_slide(prs)
    build_eda_hist_slide(prs)
    build_approach_slide(prs)
    build_feature_slide(prs)
    build_accuracy_slide(prs)
    build_shap_slide(prs)
    build_result_slide(prs)
    build_business_slide(prs)
    build_vibe_slide(prs)
    build_summary_slide(prs)

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PATH))
    return OUTPUT_PATH


def main() -> None:
    generate_eda_charts()
    output = build_presentation()
    print(f"Saved presentation: {output}")


if __name__ == "__main__":
    main()
